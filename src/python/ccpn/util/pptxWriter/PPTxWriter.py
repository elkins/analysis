"""
    python-pptx is a Python library for creating, modifying, and extracting data from PowerPoint (.pptx) files.

    Hierarchy:
    •	Presentation: Represents the whole PowerPoint file.
        •	Slide Masters: Templates that define the global layout and design for the slides.
            •	Slide Layouts: Predefined slide formats based on a slide master. These layouts are used for creating individual slides.
                    •	Shapes: Objects that can be added to a slide (e.g., text boxes, tables, images).
                        •	Placeholders: Predefined content areas on a slide layout (e.g., title, body text, image, etc.).

NB: For simplicity and practicality is often used only one single slide master in presentations, therefore also here we use only the first slide master and its layouts as template
Multiple slide masters are messy and maybe used for Complex Presentations or joint presentations...

"""

import pandas as pd
import importlib
import matplotlib.pyplot as plt
from abc import ABC, abstractmethod
from collections import defaultdict
from zipfile import ZipFile
from typing import Union, cast
from pptx import Presentation as _presentation
from pptx.slide import Slides, Slide, SlideLayout
from pptx.shapes.placeholder import LayoutPlaceholder
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR,  MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from xml.etree.ElementTree import parse, register_namespace
from ccpn.util.Path import aPath, checkFilePath
from ccpn.framework.Application import getApplication, getMainWindow, getProject, getCurrent
from ccpn.util.Logging import getLogger

PLACEHOLDER_NAME = 'placeholder_name'
PLACEHOLDER_TYPE = 'placeholder_type'
PLACEHOLDER_GETTER = 'placeholder_getter'
PLACEHOLDER_SETTER = 'placeholder_setter'
PLACEHOLDER_STYLE_GETTER = 'placeholder_style_getter'
PLACEHOLDER_STYLE_SETTER = 'placeholder_style_setter'

PLACEHOLDER_TYPE_TEXT = 'Text'
PLACEHOLDER_TYPE_IMAGE = 'Image'
PLACEHOLDER_TYPE_TABLE = 'Table'


class PPTStyleManager():

    # ccpn colours
    CCPNpurple = (106, 59, 113) #6A3B71
    CCPNgreen = (47, 112, 92) #2F705C
    CCPNyellow = (189, 157, 70) #BD9D46
    CCPNblue = (12, 79, 131) #0C4F83'

    tableStyle = {
        "header_font_size"                              : 10,
        "header_font_name"                           : "Helvetica",
        "header_font_color"                            : (255, 255, 255),  # White text color for the header
        "header_bold"                                     : True,
        "header_alignment"                            : "center",
        "header_background_color"               : CCPNpurple,
        "font_size"                                          : 10,
        "font_name"                                       : "Helvetica",
        "font_color"                                        : (0, 0, 0),  # Black text color for both rows
        "bold"                                                 : False,
        "alignment"                                        : "center",
        "border_color"                                   : (0, 0, 0),
        "border_width"                                  : 1,
        "cell_padding"                                   : 1,
        "cell_background_color"                   : (239, 235, 240),
        "cell_background_alternate_color"   : (200, 193, 201),
        }

    """Class to manage styles in PowerPoint, including placeholders and table styling."""

    def __init__(self, pptPresenation):
        """
        Initialize the style manager with a presentation instance.
        :param presentation: An instance of a PowerPoint Presentation.
        """
        self.pptPresenation = pptPresenation
        self._accentColours = self._extractThemeAccents()

    def _extractThemeAccents(self):
        """Extracts accent colors from a PowerPoint presentation's theme.xml."""
        accentColors = {}

        # Open the PowerPoint file as a zip archive
        with ZipFile(self.pptPresenation._pptxPath, 'r') as pptxZip:
            # Locate and open theme XML file
            themeFile = [f for f in pptxZip.namelist() if 'theme/theme' in f][0]
            with pptxZip.open(themeFile) as themeXml:
                themeTree = parse(themeXml)

                # Define namespaces
                namespaces = {
                    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    }
                register_namespace('', namespaces['a'])

                # Extract colors from theme
                clrScheme = themeTree.find('.//a:clrScheme', namespaces)
                if clrScheme is not None:
                    for accent in range(1, 7):  # Accent1 to Accent6
                        accentElem = clrScheme.find(f".//a:accent{accent}/a:srgbClr", namespaces)
                        if accentElem is not None:
                            rgbValue = accentElem.get('val')
                            accentColors[f'accent{accent}'] = RGBColor.from_string(rgbValue)
        return accentColors

    def getTextStyle(self, layoutPlaceholder):
        """
        Parses the XML of a layout placeholder's textFrame and extracts all text style information.
        :param layoutPlaceholder: The placeholder shape from the layout slide.
        :param accentColours: A dictionary mapping schemeClr values to RGB colors.
        :return: A dictionary containing all text styles and attributes in camelCase.
        """
        fontStyles = {}

        # Access the shape XML
        shapeXml = layoutPlaceholder._element

        # Find the <p:txBody> element in the XML
        namespaces = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
            }
        textBody = shapeXml.find('.//p:txBody', namespaces)
        if textBody is None:
            return fontStyles  # No textFrame found, return empty dict

        # Parse <a:bodyPr> for body properties
        bodyPr = textBody.find('a:bodyPr', namespaces)
        if bodyPr is not None:
            fontStyles['anchor'] = bodyPr.get('anchor', None)
            fontStyles['autofit'] = bodyPr.find('a:normAutofit', namespaces) is not None
            fontStyles['wrap'] = bodyPr.get('wrap', None)

        # Parse <a:lstStyle> for list style properties
        lstStyle = textBody.find('a:lstStyle', namespaces)
        if lstStyle is not None:
            for level in range(1, 10):  # Check levels 1 to 9
                lvlTag = f'a:lvl{level}pPr'
                lvlPPr = lstStyle.find(lvlTag, namespaces)
                if lvlPPr is not None:
                    levelKey = f'level{level}'
                    fontStyles[levelKey] = {}
                    fontStyles[levelKey]['marginLeft'] = lvlPPr.get('marL', None)
                    fontStyles[levelKey]['indent'] = lvlPPr.get('indent', None)
                    fontStyles[levelKey]['alignment'] = lvlPPr.get('algn', None)
                    fontStyles[levelKey]['bulletType'] = 'none' if lvlPPr.find('a:buNone', namespaces) else 'default'

                    # Default run properties
                    defRPr = lvlPPr.find('a:defRPr', namespaces)
                    if defRPr is not None:
                        fontSize = defRPr.get('sz')
                        if fontSize:
                            fontStyles[levelKey]['fontSize'] = int(fontSize) / 100  # Convert from centipoints to points
                        else:
                            fontStyles[levelKey]['fontSize'] = 10  # default

                        fontStyles[levelKey]['bold'] = defRPr.get('b', None)

                        # Extract font typeface
                        latin = defRPr.find('a:latin', namespaces)
                        if latin is not None:
                            fontStyles[levelKey]['fontTypeface'] = latin.get('typeface', None)

                        # Extract font color
                        solidFill = defRPr.find('a:solidFill', namespaces)
                        if solidFill is not None:
                            # Check for schemeClr
                            schemeClr = solidFill.find('a:schemeClr', namespaces)
                            if schemeClr is not None:
                                schemeColorName = schemeClr.get('val')
                                if schemeColorName in self._accentColours:
                                    fontStyles[levelKey]['fontColor'] = self._accentColours[schemeColorName]
                                else:
                                    fontStyles[levelKey]['fontColor'] = "RGB(0,0,0)"  # Default if schemeClr isn't mapped
                            else:
                                # Check for solid RGB color
                                srgbClr = solidFill.find('a:srgbClr', namespaces)
                                if srgbClr is not None:
                                    fontStyles[levelKey]['fontColor'] = f"#{srgbClr.get('val')}"

        # Parse text paragraphs and runs
        fontStyles['paragraphs'] = []
        for paragraph in textBody.findall('a:p', namespaces):
            paraDict = {}

            # Paragraph properties
            pPr = paragraph.find('a:pPr', namespaces)
            if pPr is not None:
                paraDict['level'] = pPr.get('lvl', None)
                defRPr = pPr.find('a:defRPr', namespaces)
                if defRPr is not None:
                    paraDict['defaultFontSize'] = defRPr.get('sz', 10)

            # Extract text runs
            paraDict['runs'] = []
            for run in paragraph.findall('a:r', namespaces):
                runDict = {}
                rPr = run.find('a:rPr', namespaces)
                if rPr is not None:
                    runDict['language'] = rPr.get('lang', None)
                    runDict['dirty'] = rPr.get('dirty', None)
                text = run.find('a:t', namespaces)
                if text is not None:
                    runDict['text'] = text.text
                paraDict['runs'].append(runDict)

            # Add paragraph to the paragraphs list
            fontStyles['paragraphs'].append(paraDict)

        return fontStyles

    def _insertTextAndApplyStyle(self, text, textBox, layoutPlaceholder, fontSizePt=None):
        """Insert the given text to the textBox and apply the layoutPlaceholder style.
        Adjust font size to fit the box if necessary, without resetting the size after.
        """
        textStyle = self.getTextStyle(layoutPlaceholder)
        # Get the text frame from the placeholder
        textFrame = textBox.text_frame

        # Apply body properties (don't reset font size here, leave it for later)
        if 'anchor' in textStyle:
            # Map anchor strings to MSO_VERTICAL_ANCHOR enum values
            anchor_map = {
                'top'   : MSO_VERTICAL_ANCHOR.TOP,
                'middle': MSO_VERTICAL_ANCHOR.MIDDLE,
                'bottom': MSO_VERTICAL_ANCHOR.BOTTOM
                }
            # Default to TOP if the anchor is not valid
            textFrame.vertical_anchor = anchor_map.get(textStyle['anchor'], MSO_VERTICAL_ANCHOR.TOP)

        if 'autofit' in textStyle:
            if textStyle['autofit']:
                textFrame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            else:
                textFrame.auto_size = MSO_AUTO_SIZE.NONE

        # Apply paragraph-level styles
        for i, paragraphData in enumerate(textStyle.get('paragraphs', [])):
            # Ensure the paragraph exists
            if i >= len(textFrame.paragraphs):
                paragraph = textFrame.add_paragraph()
            else:
                paragraph = textFrame.paragraphs[i]

            # Apply paragraph properties
            if 'level' in paragraphData:
                paragraph.level = int(paragraphData['level']) if paragraphData['level'] is not None else 0

            # Apply runs (text segments)
            paragraph.clear()  # Clear any existing runs
            for runData in paragraphData.get('runs', []):
                run = paragraph.add_run()
                run.text = text
                rPr = run.font

                # Apply font properties for the run
                if fontSizePt:
                    rPr.size = fontSizePt  # Use the pre-adjusted font size

                if 'language' in runData:
                    run.language_id = runData['language']
                if 'fontTypeface' in textStyle.get(f"level{paragraph.level + 1}", {}):
                    rPr.name = textStyle[f"level{paragraph.level + 1}"]['fontTypeface']
                if 'fontColor' in textStyle.get(f"level{paragraph.level + 1}", {}):
                    fontColor = textStyle[f"level{paragraph.level + 1}"]['fontColor']
                    if fontColor.startswith('#'):
                        rPr.color.rgb = RGBColor.from_string(fontColor[1:])
                    elif fontColor.startswith('RGB'):
                        rgb = fontColor[4:-1].split(',')
                        rPr.color.rgb = RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))

                # Apply other run styles
                if 'bold' in textStyle.get(f"level{paragraph.level + 1}", {}):
                    rPr.bold = textStyle[f"level{paragraph.level + 1}"]['bold'] == '1'
                break

        # Apply list level styles (e.g., margins, alignment) for each level
        for level in range(1, 10):
            levelKey = f"level{level}"
            if levelKey in textStyle:
                levelStyles = textStyle[levelKey]
                for paragraph in textFrame.paragraphs:
                    try:
                        if paragraph.level == level - 1:
                            if 'marginLeft' in levelStyles:
                                paragraph.margin_left = Inches(float(levelStyles['marginLeft']) / 914400)  # EMUs to Inches
                            if 'indent' in levelStyles:
                                paragraph.space_before = Inches(float(levelStyles['indent']) / 914400)  # EMUs to Inches
                            if 'alignment' in levelStyles:
                                alignmentMap = {
                                    'l'      : PP_ALIGN.LEFT,
                                    'ctr'    : PP_ALIGN.CENTER,
                                    'r'      : PP_ALIGN.RIGHT,
                                    'justify': PP_ALIGN.JUSTIFY
                                    }
                                paragraph.alignment = alignmentMap.get(levelStyles['alignment'], PP_ALIGN.LEFT)
                    except Exception as styleError:
                        print(f'Error applying the style from template {text}, {textBox}, {layoutPlaceholder}. Error: {styleError}')

    def _addTextToTextBox(self, text, textBox, layoutPlaceholder):
        """Adjust the font size of the text to ensure it fits in the given text box."""
        # Get text style
        textStyle = self.getTextStyle(layoutPlaceholder)
        textFrame = textBox.text_frame
        # Retrieve container height and width
        containerHeight = textBox.height
        containerWidth = textBox.width
        # Default font size (use the default or initial size from the style)
        fontSize = textStyle.get('defaultFontSize', 12)  # Assume 12pt as a default
        # Set the initial font size
        fontSizePt = Pt(fontSize)
        # Create a function to estimate the total text height (in points)
        def _estimateTextHeight(text, fontSizePt, containerWidth):
            """Estimate the total height of the text based on the font size and text length."""
            estimatedHeight = 0
            lines = text.split('\n')  # Split the text into lines
            # Estimate the height per line (using the font size)
            lineHeight = fontSizePt * 1.2  # 1.2 multiplier for line spacing (can be adjusted)
            for line in lines:
                # Estimate the width of the line based on the number of characters
                # This is a simplification, and you could improve it by taking actual font metrics into account
                numChars = len(line)
                lineWidth = numChars * fontSizePt / 2  # Approximate character width (can adjust multiplier)
                # If line width exceeds the container width, we would need to break it into more lines
                if lineWidth > containerWidth:
                    # Estimate how many lines are needed
                    numLines = (lineWidth // containerWidth) + 1
                    estimatedHeight += numLines * lineHeight
                else:
                    # If it fits in one line, just add the line height
                    estimatedHeight += lineHeight
            return estimatedHeight

        # Estimate the text height with the initial font size
        estimatedHeight = _estimateTextHeight(text, fontSizePt, containerWidth)
        # Check if estimated height exceeds the available space
        while estimatedHeight > containerHeight and fontSizePt > Pt(1):  # Ensure font doesn't go below 1pt
            fontSizePt -= Pt(1)  # Reduce font size by 1pt
            estimatedHeight = _estimateTextHeight(text, fontSizePt, containerWidth)  # Recalculate height

        # Apply the adjusted font size to the text style
        self._insertTextAndApplyStyle(text, textBox, layoutPlaceholder, fontSizePt)  # Pass the adjusted font size

    def applyTableStyle(self, table):
        # Apply style for header row (first row)
        headerRow = table.rows[0]
        styleDict = self.tableStyle
        for cell in headerRow.cells:
            cell.text_frame.paragraphs[0].font.size = Pt(styleDict["header_font_size"])
            cell.text_frame.paragraphs[0].font.name = styleDict["header_font_name"]
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(*styleDict["header_font_color"])
            cell.text_frame.paragraphs[0].font.bold = styleDict["header_bold"]
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if styleDict["header_alignment"] == "center" else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*styleDict["header_background_color"])
            # Apply padding for header cells
            cell.margin_top = styleDict["cell_padding"]
            cell.margin_bottom = styleDict["cell_padding"]
            cell.margin_left = styleDict["cell_padding"]
            cell.margin_right = styleDict["cell_padding"]
            #  border styling seems not working

        # Apply style for all rows except the header
        for rowIndex, row in enumerate(list(table.rows)[1:], start=1):
            for colIndex, cell in enumerate(row.cells):
                # cell.text_frame.word_wrap = True
                # cell.text_frame.auto_size = None  # Disable auto-size
                cell.text_frame.paragraphs[0].font.size = Pt(styleDict["font_size"])
                cell.text_frame.paragraphs[0].font.name = styleDict["font_name"]
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(*styleDict["font_color"])
                cell.text_frame.paragraphs[0].font.bold = styleDict["bold"]
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if styleDict["alignment"] == "center" else PP_ALIGN.LEFT
                # Set vertical alignment to center for cells
                cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                # Alternating row background colors
                if rowIndex % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*styleDict["cell_background_alternate_color"])
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*styleDict["cell_background_color"])

                # Apply padding for table cells
                cell.margin_top = styleDict["cell_padding"]
                cell.margin_bottom = styleDict["cell_padding"]
                cell.margin_left = styleDict["cell_padding"]
                cell.margin_right = styleDict["cell_padding"]


class PPTxPresentation():

    def __init__(self, presentationTemplate, placeholderErrorPolicy='raise'):
        """
        :param pptxPath: The path to an existing PPTx presentation containing a single slide master and layout(s) with appropriately
                                     named placeholders to use as a template for a new presentation.
        """
        self._presentationTemplate = presentationTemplate
        self._pptxPath = self._presentationTemplate.getAbsoluteTemplatePath()
        self._presentation = _presentation(self._pptxPath)
        self._placeholderErrorPolicy = placeholderErrorPolicy
        self._data = None

        if len(self._presentation.slide_masters)>1:
            raise RuntimeWarning('Multiple slide masters are not supported. This  Presentation will use only the first slide master')
        self.styleManager = PPTStyleManager(self)

    @property
    def data(self):
        return self._data

    def setData(self, **kwargs):
        self._data = {**kwargs} # This creates a shallow copy of the kwargs dictionary instead of data = kwargs (that points to the same passed in dict)

    @abstractmethod
    def buildFromTemplate(self):
        """
        Builds a new presentation based on the template, dynamically applying content to placeholders
        as defined in the slide mapping.
        """
        self._presentationTemplate.setData(**self.data)
        isValidTemplate, templateErrors = self._validateTemplate()
        if not isValidTemplate and self._placeholderErrorPolicy == 'raise':
            raise RuntimeError(f'Detected errors while building a new Presentation from Template \n{self._formatDefaultDict(templateErrors)}')
        else:

            slideMapping = self._presentationTemplate.slideMapping
            for slideLayoutName, placeholderDefs in slideMapping.items():
                layout = self.getLayout(slideLayoutName)
                newSlide = self.newSlide(layout, removePlaceholders=True)
                for placeholderDef in placeholderDefs:
                    try:
                        self._handlePlaceholder(newSlide, layout, placeholderDef)
                    except Exception as ex:
                        print(f'Some Error in filling the placeholder occurred: {ex}')

    def save(self, filePath):
        self._presentation.save(filePath)

    def newSlide(self, layoutIdentifier: Union[SlideLayout, int, str], removePlaceholders: bool=True) -> Slide:
        """
        A method to add slide and clone (if necessary) the placeholder from the master slide .
        Defaults remove all newly generated placeholders. This is intentional because there are too many issues/bugs in the release of pptx used when creating this class.
        Furthermore, placeholders   in the new slides are confusing with the master layout placeholders, because not exactly as the template.
        Pptx when cloning the placeholders gives new names and shape ids and makes impossible to back-track the original placeholder names, defeating
        the whole reason naming placeholders in the master template! The only way to find it back seems by matching by placeholder.element.ph_idx [ Took a while to spot this! ]
        If we keep the placeholders, We need to rename them to match the placeholder from the master slide.
        :return: a new slide
        """
        if layoutIdentifier is None:
            raise ValueError("A layout must be provided.")
        if isinstance(layoutIdentifier, SlideLayout):
            layout = layoutIdentifier
        elif isinstance(layoutIdentifier, (int, str)):
            layout = self.getLayout(layoutIdentifier)
        else:
            raise ValueError("The layoutIdentifier parameter must be a SlideLayout, int, or str.")

        newSlide = self._presentation.slides.add_slide(layout)
        if removePlaceholders:
            self._removePlaceholdersForSlide(newSlide)
        else:
            self._matchPlaceholdersNameToLayout(newSlide, layout)
        return newSlide

    def getSlides(self) -> tuple[Slide, ...]:
        """A list of slides of this presentation."""
        sldIdLst = self._presentation._element.get_or_add_sldIdLst()
        self._presentation.part.rename_slide_parts([cast("CT_SlideId", sldId).rId for sldId in sldIdLst])
        return tuple([slide for slide in Slides(sldIdLst, self._presentation)])

    def getLayout(self, identifier: int | str = 0) -> SlideLayout:
        """
        Retrieve a slide layout by index or name.

        :param identifier: Either an integer (layout index) or a string (layout name).
        :return: The slide layout object.
        :raises ValueError: If the identifier is invalid or out of range.
        :raises RuntimeError: If a name is provided but no matching layout is found.
        """
        layouts = self._presentation.slide_layouts

        # Handle identifier as an integer (index)
        if isinstance(identifier, int):
            if 0 <= identifier < len(layouts):
                return layouts[identifier]
            raise ValueError(f"Layout index {identifier} is out of range. Valid indices: 0-{len(layouts) - 1}")

        # Handle identifier as a string (name)
        elif isinstance(identifier, str):
            for layout in layouts:
                if layout.name == identifier:
                    return layout
            raise RuntimeError(f"Could not find a layout called: {identifier}")

        # Invalid identifier type
        raise ValueError("Identifier must be an integer (index) or a string (name).")

    def getLayoutPlaceholders(self, layoutIdentifier: Union[SlideLayout, int, str]) -> dict[str, LayoutPlaceholder]:
        """
        Retrieve the placeholders from a slide layout.

        :param layoutIdentifier: The slide layout, identified by one of the following:
                                  - SlideLayout object
                                  - Integer (layout index)
                                  - String (layout name)
        :return: A dictionary of placeholders in the specified layout.
        """
        if layoutIdentifier is None:
            raise ValueError("A layout must be provided.")
        if isinstance(layoutIdentifier, SlideLayout):
            slideLayout = layoutIdentifier
        elif isinstance(layoutIdentifier, (int, str)):
            slideLayout = self.getLayout(layoutIdentifier)
        else:
            raise ValueError("The layoutIdentifier parameter must be a SlideLayout, int, or str.")
        return {placeholder.name: placeholder for placeholder in slideLayout.placeholders if placeholder.is_placeholder and placeholder.name}

    def insertText(self, targetSlide, templateLayout, placeholderName, text):
        """Inserts text into a PowerPoint slide using an existing placeholder shape."""
        layoutPlaceholder = self._getLayoutPlaceholder(placeholderName, templateLayout)
        left, top, width, height = self._getLayoutPlaceholderCoords(layoutPlaceholder)
        textBoxShape = targetSlide.shapes.add_textbox(left, top, width, height)
        self.styleManager._addTextToTextBox(text, textBoxShape, layoutPlaceholder)
        return textBoxShape

    def insertDataFrame(self, targetSlide, templateLayout, placeholderName, dataFrame, fontSize=7):
        """Converts a Pandas DataFrame to a PowerPoint table using an existing shape.
        """
        layoutPlaceholder = self._getLayoutPlaceholder(placeholderName, templateLayout)
        colnames = list(dataFrame.columns)
        rows, cols = dataFrame.shape

        # Get position and size from placeholder
        left, top, width, height = self._getLayoutPlaceholderCoords(layoutPlaceholder)

        # Create the table in the slide with same position and size as the placeholder
        tableShape = targetSlide.shapes.add_table(rows + 1, cols, left, top, width, height)  # +1 for header row
        table = tableShape.table

        # Fill in the column headers (first row)
        for colIndex, colName in enumerate(colnames):
            cell = table.cell(0, colIndex)
            cell.text = colName
            # Make headers bold (optional)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

        # Fill in the DataFrame values (starting from the second row)
        for rowIndex, (index, row) in enumerate(dataFrame.iterrows(), start=1):  # start=1 to begin after header row
            for colIndex, val in enumerate(row):
                table.cell(rowIndex, colIndex).text = str(val)
        # Set style for all cells
        self.styleManager.applyTableStyle(table)

    def insertImage(self, targetSlide, templateLayout, placeholderName, imagePath):

        layoutPlaceholder = self._getLayoutPlaceholder(placeholderName, templateLayout)
        # Get master placeholder's dimensions
        left, top, width, height = self._getLayoutPlaceholderCoords(layoutPlaceholder)

        imageWidth, imageHeight = self._getImageSize(imagePath)
        imageHeight = imageHeight or 300 #px default fallback
        imageWidth = imageWidth or 300
        # Calculate the aspect ratio of the image
        imageRatio = imageWidth / imageHeight
        placeholderRatio = width / height

        # Scale the image to fit within the placeholder
        if imageRatio > placeholderRatio:
            # Fit by width
            scaledWidth = width
            scaledHeight = int(width / imageRatio)
        else:
            # Fit by height
            scaledHeight = height
            scaledWidth = int(height * imageRatio)

        # Center the scaled image within the placeholder
        centeredLeft = left + (width - scaledWidth) // 2
        centeredTop = top + (height - scaledHeight) // 2

        # Add the scaled and centered image to the slide
        targetSlide.shapes.add_picture(
                imagePath,
                centeredLeft,
                centeredTop,
                width=scaledWidth,
                height=scaledHeight,
                )

    @staticmethod
    def _getImageSize(filePath):
        try:
            image = plt.imread(filePath)
            height, width = image.shape[:2]  # First two dimensions are height and width
            return width, height
        except Exception as err:
            getLogger().debug(f'Cannot get image size for {filePath}. Error: {err}')
            return None, None

    @staticmethod
    def _formatDefaultDict(d):
        items = "\n".join(f"  {repr(k)}: {repr(v)}," for k, v in d.items())
        return f"{{\n{items}\n}}"

    def _validateTemplate(self):
        """
        Examine the template SlideMapping and ensure names are properly defined and any getter/setter exists
        """
        errorMsgDict = defaultdict(list)
        slideMapping = self._presentationTemplate.slideMapping
        for slideLayoutName, placeholderDefs in slideMapping.items():
            layout = None
            try:
                layout = self.getLayout(slideLayoutName)
            except ValueError as err:
                errorMsgDict[slideLayoutName].append(err)
            if layout is not None:
                for placeholderDef in placeholderDefs:
                    placeholderName = placeholderDef.get(PLACEHOLDER_NAME)
                    placeholderType = placeholderDef.get(PLACEHOLDER_TYPE)
                    placeholderGetter = placeholderDef.get(PLACEHOLDER_GETTER)
                    ph = None
                    try:
                        ph = self._getLayoutPlaceholder(placeholderName, layout)
                    except ValueError as phErr:
                        errorMsgDict[(slideLayoutName, placeholderName)].append(phErr)

                    if ph:
                        getterFunc = getattr(self._presentationTemplate, placeholderGetter, None)
                        if getterFunc is None:
                            errorMsgDict[(slideLayoutName, placeholderName)].append(f'Could not find a valid getter: {placeholderGetter} is not defined in the template Class')
                        #  we have getter  in the class . now er need to check  the value returned is the same type from the one defined in the mapping
                        else:
                            try:
                                value = getterFunc()
                                if isinstance(value, str) and placeholderType in [PLACEHOLDER_TYPE_TEXT, PLACEHOLDER_TYPE_IMAGE]:
                                    continue
                                elif isinstance(value, pd.DataFrame) and placeholderType in [PLACEHOLDER_TYPE_TABLE]:
                                    continue
                                else:
                                    errorMsgDict[(slideLayoutName, placeholderName)].append(f'The placeholder Type is not same as getter')
                            except Exception as _ex:
                                errorMsgDict[(slideLayoutName, placeholderName)].append(f'Error retrieving the placeholder value. {_ex}')

        isValid = True if len(errorMsgDict) == 0 else False
        return isValid, errorMsgDict


    def _getLayouts(self):
        """
        Get the master slide layouts. Layouts contain the master placeholders used to create any new slides
        :return: list of slide layouts
        """
        return [slide for slide in self._presentation.slide_layouts]

    def slideWidth(self):
        return self._presentation.slide_width

    def setSlideWidth(self, width):
        self._presentation.slide_width = width

    def slideHeight(self):
        return self._presentation.slide_height

    def setSlideHeight(self, height):
        self._presentation.slide_height = height

    def _matchPlaceholdersNameToLayout(self, slide, layout):
        for newPlaceholder in slide.placeholders:
            for masterPlaceholder in layout.placeholders:
                if masterPlaceholder.element.ph_idx == newPlaceholder.element.ph_idx:
                    newPlaceholder.name = masterPlaceholder.name

    def _removePlaceholdersForSlide(self, slide):
        for shape in list(slide.shapes):  # Convert to list to avoid modification issues
            if shape.is_placeholder:
                shapePlaceholderFormat = shape.placeholder_format
                shapePlaceholderFormat.element.getparent().remove(shapePlaceholderFormat.element)
        return slide

    def _getLayoutPlaceholder(self, placeholderName, templateLayout) -> LayoutPlaceholder:
        layoutPlaceholders = self.getLayoutPlaceholders(templateLayout)
        layoutPlaceholder = layoutPlaceholders.get(placeholderName, None)
        if layoutPlaceholder is None:
            raise ValueError(f'Could not find the placeholder called {placeholderName}')
        return layoutPlaceholder

    def _getLayoutPlaceholderCoords(self, placeholder):
        left = placeholder.left
        top = placeholder.top
        width = placeholder.width
        height = placeholder.height
        return (left, top, width, height)

    def _handlePlaceholder(self, slide, layout, placeholderDef,  **getterKwargs):
        """
        Handles insertion of content into a placeholder based on its type.
        """
        placeholderName = placeholderDef.get(PLACEHOLDER_NAME)
        placeholderType = placeholderDef.get(PLACEHOLDER_TYPE)
        placeholderGetter = placeholderDef.get(PLACEHOLDER_GETTER)
        # Dynamically retrieve the content for the placeholder
        getterFunc = getattr(self._presentationTemplate, placeholderGetter, None)
        if getterFunc is None:
            return
        value = getterFunc(**getterKwargs)
        if placeholderType == 'Text':
            self.insertText(slide, layout, placeholderName, value or '')
        elif placeholderType == PLACEHOLDER_TYPE_IMAGE:
            isPathOk, msgErr = checkFilePath(aPath(value))
            if isPathOk:
                self.insertImage(slide, layout, placeholderName, value)
        elif placeholderType == PLACEHOLDER_TYPE_TABLE:
            try:
                self.insertDataFrame(slide, layout, placeholderName, value)
            except Exception as err:
                print('Cannot add table', err)


class ScreenPPTxPresentation(PPTxPresentation):

    def buildFromTemplate(self):
        import ccpn.AnalysisScreen.lib.experimentAnalysis.matching.MatchingVariables as mv
        self._presentationTemplate.setData(**self.data)
        isValidTemplate, templateErrors = self._validateTemplate()
        if not isValidTemplate and self._placeholderErrorPolicy == 'raise':
            raise RuntimeError(f'Detected errors while building a new Presentation from Template \n{self._formatDefaultDict(templateErrors)}')
        else:
            slideMapping = self._presentationTemplate.slideMapping
            # build the title (first slide)
            if not slideMapping:
                raise RuntimeError(f'Detected errors while building a new Presentation from Template \n{self._formatDefaultDict(templateErrors)}')
            titleLayoutName = list(slideMapping.keys())[0]
            titlePlaceholderDefs = slideMapping[titleLayoutName]
            layout = self.getLayout(titleLayoutName)
            newSlide = self.newSlide(layout, removePlaceholders=True)
            for placeholderDef in titlePlaceholderDefs:
                try:
                    self._handlePlaceholder(newSlide, layout, placeholderDef)
                except Exception as ex:
                    print(f'Some Error in filling the placeholder occurred: {ex}')

            # build the substances Slides
            substanceTable = self.data.get('substanceTable')
            matchingTable = self.data.get('matchingTable')
            if substanceTable is None:
                return

            for tableIndex, substanceTableRow in substanceTable.iterrows():
                substancePid = substanceTableRow[mv.Reference_SubstancePid]
                matchingTableForSubstance = matchingTable[matchingTable[mv.Reference_SubstancePid]==substancePid]
                titleLayoutName = list(slideMapping.keys())[1]
                titlePlaceholderDefs = slideMapping[titleLayoutName]
                layout = self.getLayout(titleLayoutName)
                newSlide = self.newSlide(layout, removePlaceholders=True)
                for placeholderDef in titlePlaceholderDefs:
                    self._handlePlaceholder(newSlide, layout, placeholderDef,
                                            substanceTableRow=substanceTableRow,
                                            matchingTableForSubstance=matchingTableForSubstance)



class PresentationTemplateABC(ABC):

    templateName = 'PresentationTemplateABC' # the name that will appear in the GUI selections
    templateEntryOrder = -1 # the order in which will appear in the GUI selections
    templateRelativePath = '' # the relative path from this file where the template is located
    slideMapping = {}

    def __init__(self, *args, **kwargs):
        self.project = getProject()
        self.application = getApplication()
        self.mainWindow = getMainWindow()
        self.current = getCurrent()
        self._data = None

    @property
    def data(self):
        return self._data

    def setData(self, **kwargs):
        self._data = {**kwargs}

    def getAbsoluteTemplatePath(self):
        # Get the module where the current class (or subclass) is defined
        moduleName = self.__class__.__module__
        module = importlib.import_module(moduleName)
        # Retrieve the file path of  the subclassed module
        thisFile = aPath(module.__file__)
        workingDir = aPath(thisFile.parent)
        absPath = workingDir / self.templateRelativePath
        return absPath

    @staticmethod
    def formatNestedDictToText(data, indentLevel=0):
        """
        Converts a nested dictionary into a readable indented string.
        :param data: The nested dictionary.
        :param indentLevel: The current level of indentation.
        :return: A formatted string.
        """
        formattedText = ""
        indent = "  " * indentLevel  # Two spaces per level
        for key, value in data.items():
            if isinstance(value, dict):
                formattedText += f"{indent}{key}:\n"
                formattedText += PresentationTemplateABC.formatNestedDictToText(value, indentLevel + 1)
            else:
                formattedText += f"{indent}{key}: {value}\n"
        return formattedText